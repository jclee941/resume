from pptx.util import Pt

from .pptx_engine import get_current_age
from .pptx_layouts import PPTXLayoutProfile
from .pptx_utils import (
    hide_empty_rows,
    join_truncate,
    resize_table_text,
    set_cell_text,
    truncate,
)


def handle_ta_profile(tbl, source, profile: PPTXLayoutProfile):
    layout = profile.table_layout(15, 8)
    resize_table_text(
        tbl,
        header_rows=layout.header_rows,
        header_cols=layout.header_columns,
        header_size_pt=layout.header_size_pt,
        body_size_pt=layout.body_size_pt,
    )
    if tbl.cell(0, 0).text_frame.text.strip():
        for paragraph in tbl.cell(0, 0).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(profile.title_size_pt)

    personal = source["personal"]
    education = source["education"]
    summary = source["summary"]
    current = source["current"]
    latest_career = source["careers"][0] if source["careers"] else {}
    set_cell_text(tbl.cell(2, 0), personal["name"])
    set_cell_text(tbl.cell(2, 1), current.get("company", latest_career.get("company", "")).replace("(주)", ""))
    set_cell_text(tbl.cell(2, 2), current.get("position", latest_career.get("role", "")))
    set_cell_text(tbl.cell(2, 4), summary.get("grade", latest_career.get("role", "")))
    set_cell_text(tbl.cell(2, 5), f"{education['school']}({education['status']})")
    set_cell_text(tbl.cell(2, 6), education["major"])
    birth_date = personal.get("birthDate")
    age = f"만   {get_current_age(birth_date)} 세" if birth_date else ""
    set_cell_text(tbl.cell(3, 1), age)
    set_cell_text(tbl.cell(3, 4), summary["totalExperience"])
    set_cell_text(tbl.cell(2, 7), ", ".join(summary.get("expertise", [])))

    rows = len(tbl.rows)
    for index, career in enumerate(source["careers"][: layout.max_entries]):
        row_index = 8 + index
        if row_index < rows:
            set_cell_text(tbl.cell(row_index, 0), truncate(career["company"], "company"))
            set_cell_text(tbl.cell(row_index, 1), career["period"])
            set_cell_text(tbl.cell(row_index, 4), truncate(career["description"], "description"))
    for index, certification in enumerate(source.get("certifications", [])[: layout.max_entries]):
        row_index = 8 + index
        if row_index < rows:
            set_cell_text(tbl.cell(row_index, 5), certification["name"])
            set_cell_text(tbl.cell(row_index, 7), certification.get("date", ""))


def handle_ta_projects(tbl, source, profile: PPTXLayoutProfile):
    layout = profile.table_layout(12, 6)
    resize_table_text(
        tbl,
        header_rows=layout.header_rows,
        header_size_pt=layout.header_size_pt,
        body_size_pt=layout.body_size_pt,
    )
    if tbl.cell(0, 0).text_frame.text.strip():
        for paragraph in tbl.cell(0, 0).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(profile.title_size_pt)

    rows = len(tbl.rows)
    for index, project in enumerate(source["projects"][: layout.max_entries]):
        row_index = 2 + index
        if row_index < rows:
            set_cell_text(tbl.cell(row_index, 0), truncate(project["name"], "project_name"))
            set_cell_text(tbl.cell(row_index, 1), project["role"])
            set_cell_text(tbl.cell(row_index, 2), project.get("description", ""))
            set_cell_text(
                tbl.cell(row_index, 3),
                join_truncate(
                    project["technologies"],
                    "technologies",
                    max_items=layout.max_technology_items,
                ),
            )
            set_cell_text(tbl.cell(row_index, 4), project["period"].replace(" ~ ", "~"))
            set_cell_text(tbl.cell(row_index, 5), project["client"])
    if layout.hide_empty_rows:
        hide_empty_rows(tbl, start_row=2, check_col=0)
