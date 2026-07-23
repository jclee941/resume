#!/usr/bin/env python3
"""
PPTX Generation Engine - Common logic extracted from generate_ta/shinhan_pptx.py

Usage:
    from pptx_engine import generate
    from pptx_templates import TEMPLATES
    generate(TEMPLATES["ta"], source_data)
"""

import datetime
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, TypeAlias

try:
    from pptx import Presentation
except ImportError:
    import sys
    print("python-pptx required: pip install python-pptx")
    sys.exit(1)

from pptx.util import Pt
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table
from .pptx_layouts import (
    PPTXLayoutProfile,
    SlideSizeMismatchError,
)
from .pptx_publication import publish_pptx
from .pptx_utils import apply_korean_font_to_table

logging.basicConfig(level=logging.INFO, format="%(message)s")
logger = logging.getLogger(__name__)

JSONValue: TypeAlias = str | int | float | bool | None | list["JSONValue"] | dict[str, "JSONValue"]
ResumeData: TypeAlias = dict[str, JSONValue]
TableHandler: TypeAlias = Callable[[Table, ResumeData, PPTXLayoutProfile], None]

@dataclass(frozen=True, slots=True)
class TemplateSpec:
    """Template specification for PPTX generation."""
    name: str
    source_path: Path
    template_path: Path
    output_path: Path
    layout_profile: PPTXLayoutProfile
    # Maps (rows, cols) -> handler function
    table_handlers: dict[tuple[int, int], TableHandler] = field(default_factory=dict)


def normalize_slide_title(slide, size_pt: int = 24) -> bool:
    """
    Normalize slide title to specified font size.
    Returns True if title was found and normalized.
    
    Strategy:
    1. Use slide.shapes.title if available
    2. Fallback: find topmost text shape within 100pt from top
    """
    # Try official title shape first
    if slide.shapes.title:
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size_pt)
        return True
    
    # Fallback: find topmost text shape
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if text_shapes:
        topmost = min(text_shapes, key=lambda s: s.top)
        if topmost.top < Pt(100):
            for p in topmost.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size_pt)
            return True
    
    return False


def get_current_age(birth_date_str: str, today: datetime.date | None = None) -> int:
    """Calculate current age from birth date string (YYYY.MM.DD or YYYY-MM-DD)."""
    birth_date = datetime.date.fromisoformat(birth_date_str.replace(".", "-"))
    current_date = today or datetime.date.today()
    birthday_has_not_arrived = (current_date.month, current_date.day) < (
        birth_date.month,
        birth_date.day,
    )
    return current_date.year - birth_date.year - int(birthday_has_not_arrived)


def generate(spec: TemplateSpec, source: ResumeData | None = None) -> Path:
    """
    Generate PPTX from template using source data.
    
    Args:
        spec: TemplateSpec with paths and table handlers
        source: Optional source data dict. If None, loads from spec.source_path
    
    Returns:
        Path to generated output file
    """
    # Load source if not provided
    if source is None:
        logger.info(f"📄 Loading: {spec.source_path}")
        with open(spec.source_path, "r", encoding="utf-8") as f:
            source_data: ResumeData = json.load(f)
    else:
        source_data = source
    
    logger.info(f"📝 Filling: {spec.template_path}")
    prs = Presentation(str(spec.template_path))
    actual_width_emu = int(prs.slide_width or 0)
    actual_height_emu = int(prs.slide_height or 0)
    if (actual_width_emu, actual_height_emu) != (
        spec.layout_profile.slide_width_emu,
        spec.layout_profile.slide_height_emu,
    ):
        raise SlideSizeMismatchError(
            profile_name=spec.layout_profile.name,
            expected_width_emu=spec.layout_profile.slide_width_emu,
            expected_height_emu=spec.layout_profile.slide_height_emu,
            actual_width_emu=actual_width_emu,
            actual_height_emu=actual_height_emu,
        )
    
    unhandled_tables: list[str] = []
    
    for slide_idx, slide in enumerate(prs.slides):
        # Normalize title
        normalize_slide_title(slide, spec.layout_profile.title_size_pt)
        
        # Process tables
        for shape in slide.shapes:
            if not isinstance(shape, GraphicFrame) or not shape.has_table:
                continue
            
            tbl = shape.table
            rows, cols = len(tbl.rows), len(tbl.columns)
            key = (rows, cols)
            
            handler = spec.table_handlers.get(key)
            if handler:
                handler(tbl, source_data, spec.layout_profile)
                apply_korean_font_to_table(tbl)
            else:
                unhandled_tables.append(f"Slide {slide_idx + 1}: {rows}x{cols}")
    
    # Warn about unhandled tables
    if unhandled_tables:
        logger.warning(f"⚠️  Unhandled tables: {', '.join(unhandled_tables)}")
    
    publish_pptx(prs, spec.output_path)
    
    logger.info(f"✅ Generated: {spec.output_path}")
    print(f"Layout profile: {spec.layout_profile.name}")
    return spec.output_path
